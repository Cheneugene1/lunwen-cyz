"""
文件解析模块
支持：.docx / .doc（二进制OLE格式）/ .pdf / .pptx / .txt / .xlsx / .csv / 常见图片
超长文档策略：
  - relevance_extract：按关键词相关度打分，保留高分块至字符上限
  - truncate_head：截断取前 N 字符
"""

import csv
import io
import logging
import re
import struct
from pathlib import Path
from typing import List, Optional

from .config import get
from .models import DocumentBlock, DocumentBundle

logger = logging.getLogger(__name__)

# 配置常量（从 config 读取，解析时实时调用）
def _max_chars() -> int:
    return int(get("max_parsed_chars_per_file", 400000))

def _max_pages() -> int:
    return int(get("max_document_pages_soft", 100))

def _strategy() -> str:
    return get("long_document_strategy", "relevance_extract")


# ── 各格式解析器 ──────────────────────────────────────────────

def _parse_txt(path: Path) -> List[DocumentBlock]:
    text = path.read_text(encoding="utf-8", errors="replace")
    blocks = []
    for i, para in enumerate(text.split("\n\n")):
        para = para.strip()
        if para:
            blocks.append(DocumentBlock(
                source_file=path.name, type="paragraph", text=para
            ))
    return blocks


def _parse_doc_binary(path: Path) -> List[DocumentBlock]:
    """
    读取旧版二进制 .doc（OLE2 Compound Document）文件。
    使用 olefile 读取 WordDocument 流，提取 UTF-16LE 编码的中英文文本。
    """
    try:
        import olefile
    except ImportError:
        logger.warning(".doc 文件解析需要 olefile 库：pip install olefile")
        return []

    try:
        ole = olefile.OleFileIO(str(path))
        if not ole.exists("WordDocument"):
            logger.warning("%s: 找不到 WordDocument 流", path.name)
            return []

        data = ole.openstream("WordDocument").read()
        chars = []
        i = 0
        while i < len(data) - 1:
            val = struct.unpack_from("<H", data, i)[0]
            # 保留可打印字符（含中文区间）
            if 0x20 <= val <= 0xFFFD and val != 0xFFFF:
                try:
                    ch = chr(val)
                    if ch.isprintable() or ch in "\n\r\t ":
                        chars.append(ch)
                except (ValueError, OverflowError):
                    pass
            i += 2

        raw = "".join(chars)
        # 提取有意义的段落（长度 ≥ 6 的中英文句段）
        seg_pattern = re.compile(
            r"[\u4e00-\u9fff\uff00-\uffef\w\s\u3000-\u9fff"
            r"，。；：？！、《》【】（）\[\]±°℃\-\+\=\/\.]{6,}"
        )
        segments = seg_pattern.findall(raw)

        # 合并成段落块（连续短段合并）
        blocks: List[DocumentBlock] = []
        buf = []
        for seg in segments:
            seg = seg.strip()
            if not seg:
                continue
            buf.append(seg)
            merged = " ".join(buf)
            if len(merged) > 200:
                blocks.append(DocumentBlock(
                    source_file=path.name, type="paragraph", text=merged
                ))
                buf = []

        if buf:
            blocks.append(DocumentBlock(
                source_file=path.name, type="paragraph", text=" ".join(buf)
            ))

        logger.info("_parse_doc_binary: %s → %d 段落", path.name, len(blocks))
        return blocks

    except Exception as e:
        logger.warning("_parse_doc_binary 失败 %s: %s", path.name, e)
        return []


def _parse_docx(path: Path) -> List[DocumentBlock]:
    try:
        from docx import Document  # python-docx
        doc = Document(str(path))
        blocks = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                blocks.append(DocumentBlock(
                    source_file=path.name, type="paragraph", text=text
                ))
        # 表格内容
        for table in doc.tables:
            rows = [" | ".join(cell.text.strip() for cell in row.cells)
                    for row in table.rows]
            table_text = "\n".join(rows)
            if table_text.strip():
                blocks.append(DocumentBlock(
                    source_file=path.name, type="table", text=table_text
                ))
        return blocks
    except Exception as e:
        logger.warning("docx 解析失败 %s: %s", path.name, e)
        return []


def _parse_pdf(path: Path) -> tuple[List[DocumentBlock], int]:
    """返回 (blocks, page_count)"""
    try:
        import pdfplumber
        blocks = []
        page_count = 0
        with pdfplumber.open(str(path)) as pdf:
            page_count = len(pdf.pages)
            for page in pdf.pages:
                text = page.extract_text() or ""
                for para in text.split("\n\n"):
                    para = para.strip()
                    if para:
                        blocks.append(DocumentBlock(
                            source_file=path.name,
                            type="paragraph",
                            text=para,
                            page=page.page_number,
                        ))
        return blocks, page_count
    except Exception as e:
        logger.warning("pdf 解析失败 %s: %s", path.name, e)
        return [], 0


def _parse_pptx(path: Path) -> List[DocumentBlock]:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        blocks = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = "\n".join(
                        p.text for p in shape.text_frame.paragraphs
                    ).strip()
                    if text:
                        blocks.append(DocumentBlock(
                            source_file=path.name,
                            type="paragraph",
                            text=text,
                            page=slide_idx,
                        ))
        return blocks
    except Exception as e:
        logger.warning("pptx 解析失败 %s: %s", path.name, e)
        return []


def _parse_xlsx(path: Path) -> List[DocumentBlock]:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        blocks = []
        for ws in wb.worksheets:
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                rows.append(" | ".join(cells))
            text = "\n".join(rows).strip()
            if text:
                blocks.append(DocumentBlock(
                    source_file=path.name, type="table",
                    text=f"[工作表: {ws.title}]\n{text}"
                ))
        return blocks
    except Exception as e:
        logger.warning("xlsx 解析失败 %s: %s", path.name, e)
        return []


def _parse_csv_plain(path: Path) -> List[DocumentBlock]:
    try:
        rows = []
        with open(path, encoding="utf-8", errors="replace", newline="") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(" | ".join(row))
        text = "\n".join(rows).strip()
        if text:
            return [DocumentBlock(source_file=path.name, type="table", text=text)]
        return []
    except Exception as e:
        logger.warning("csv 解析失败 %s: %s", path.name, e)
        return []


def _parse_image(path: Path) -> List[DocumentBlock]:
    """图片仅记录文件名（MVP 不做 OCR）"""
    return [DocumentBlock(
        source_file=path.name,
        type="image_caption",
        text=f"[图片文件: {path.name}，暂不支持 OCR 解析]"
    )]


# ── 超长文档策略 ──────────────────────────────────────────────

def _simple_relevance_score(block: DocumentBlock, keywords: List[str]) -> float:
    """用关键词出现次数对段落打分"""
    text_lower = block.text.lower()
    score = sum(text_lower.count(kw.lower()) for kw in keywords)
    return float(score)


def _extract_temp_keywords(blocks: List[DocumentBlock], file_name: str) -> List[str]:
    """在尚无 WritingPlan 时，从文首若干块提取临时关键词（简单高频词）"""
    # 取前 20 块文字
    sample_text = " ".join(b.text for b in blocks[:20])
    # 简单分词：中文字符连续段 + 英文单词
    words = re.findall(r"[\u4e00-\u9fff]{2,8}|[a-zA-Z]{4,}", sample_text)
    freq: dict[str, int] = {}
    for w in words:
        w_lower = w.lower()
        freq[w_lower] = freq.get(w_lower, 0) + 1
    # 去停用词（极简）
    stop = {"the", "and", "for", "that", "with", "this", "have", "from",
            "are", "were", "been", "has", "not", "but", "they", "which"}
    freq = {w: c for w, c in freq.items() if w not in stop and len(w) > 1}
    top = sorted(freq.items(), key=lambda x: -x[1])[:20]
    return [w for w, _ in top]


def _apply_long_strategy(
    blocks: List[DocumentBlock],
    file_name: str,
    strategy: str,
    max_chars: int,
) -> tuple[List[DocumentBlock], bool]:
    """
    对超长文档应用策略，返回 (精简后的 blocks, truncation_applied)
    """
    total_chars = sum(len(b.text) for b in blocks)
    if total_chars <= max_chars:
        return blocks, False

    if strategy == "truncate_head":
        # 简单截断：从头开始累加到上限
        kept = []
        acc = 0
        for b in blocks:
            if acc + len(b.text) > max_chars:
                break
            kept.append(b)
            acc += len(b.text)
        return kept, True

    # relevance_extract（默认）
    keywords = _extract_temp_keywords(blocks, file_name)
    scored = [(b, _simple_relevance_score(b, keywords)) for b in blocks]
    scored.sort(key=lambda x: -x[1])

    kept = []
    acc = 0
    for b, _ in scored:
        if acc + len(b.text) > max_chars:
            continue
        kept.append(b)
        acc += len(b.text)

    # 恢复原始顺序（按 blocks 列表中的出现位置）
    order = {id(b): i for i, b in enumerate(blocks)}
    kept.sort(key=lambda b: order.get(id(b), 999999))
    return kept, True


# ── 主入口 ────────────────────────────────────────────────────

SUPPORTED_SUFFIXES = {
    ".txt", ".docx", ".doc", ".pdf", ".pptx", ".xlsx", ".csv",
    ".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp",
}


def parse_files(
    file_paths: List[str | Path],
    keywords: Optional[List[str]] = None,
) -> DocumentBundle:
    """
    解析多个文件，返回合并后的 DocumentBundle。
    keywords 可传入已知主题词（用于 relevance_extract），否则自动提取。
    """
    all_blocks: List[DocumentBlock] = []
    is_long = False
    max_chars = _max_chars()
    max_pages = _max_pages()
    strategy = _strategy()

    for fp in file_paths:
        path = Path(fp)
        suffix = path.suffix.lower()
        page_count = 0

        if not path.exists():
            logger.warning("文件不存在，跳过：%s", path)
            continue

        # 按格式解析
        if suffix == ".txt":
            blocks = _parse_txt(path)
        elif suffix == ".docx":
            blocks = _parse_docx(path)
        elif suffix == ".doc":
            # 尝试作为 ooxml zip（部分 .doc 实质是 docx），否则走 OLE2 二进制路径
            import zipfile
            if zipfile.is_zipfile(str(path)):
                blocks = _parse_docx(path)   # python-docx 处理 ooxml
                if not blocks:
                    blocks = _parse_doc_binary(path)
            else:
                blocks = _parse_doc_binary(path)
        elif suffix == ".pdf":
            blocks, page_count = _parse_pdf(path)
        elif suffix == ".pptx":
            blocks = _parse_pptx(path)
        elif suffix == ".xlsx":
            blocks = _parse_xlsx(path)
        elif suffix == ".csv":
            blocks = _parse_csv_plain(path)
        elif suffix in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}:
            blocks = _parse_image(path)
        else:
            logger.warning("不支持的文件格式，跳过：%s", suffix)
            continue

        # 超长页数检测
        if page_count > max_pages:
            is_long = True
            logger.info("文件 %s 共 %d 页，超出软上限 %d", path.name, page_count, max_pages)

        # 超长字符检测 & 策略应用
        file_chars = sum(len(b.text) for b in blocks)
        if is_long or file_chars > max_chars:
            kws = keywords or _extract_temp_keywords(blocks, path.name)
            blocks, truncated = _apply_long_strategy(blocks, path.name, strategy, max_chars)
            if truncated:
                is_long = True
                logger.info("文件 %s 触发超长策略 (%s)，截断后 %d 块",
                            path.name, strategy, len(blocks))

        all_blocks.extend(blocks)

    # 构建 flags
    truncation_applied = is_long
    flags = {
        "long_document": is_long,
        "truncation_applied": truncation_applied,
    }
    if truncation_applied:
        flags["truncation_message"] = (
            "文档超过长度上限，已按策略仅采用部分内容参与后续规划；"
            "全文未全部进入模型。"
        )

    return DocumentBundle(blocks=all_blocks, flags=flags)
